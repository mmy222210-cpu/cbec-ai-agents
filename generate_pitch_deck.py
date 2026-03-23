import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Common styles
def set_title(slide, title_text):
    title = slide.shapes.title
    title.text = title_text
    for p in title.text_frame.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = 'Microsoft YaHei'
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 51, 102)

def set_content(slide, body_text):
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    # Split the body_text string correctly respecting indents
    for line in body_text.strip().split('\n'):
        p = tf.add_paragraph()
        if line.startswith('  -'):
            p.text = line.replace('  -', '•', 1).strip()
            p.level = 1
        elif line.startswith('-'):
            p.text = line.replace('-', '', 1).strip()
            p.level = 0
        else:
            p.text = line.strip()
            
        for run in p.runs:
            run.font.name = 'Microsoft YaHei'
            run.font.size = Pt(20)

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0] # Title layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "CBEC Agent Matrix: \n多智能体出海全链路战略引擎"
subtitle.text = "重塑中国制造的端到端“数字高管团队”\n\nAI创客大赛参赛作品"

for p in title.text_frame.paragraphs:
    for run in p.runs:
        run.font.name = 'Microsoft YaHei'
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 51, 102)

for p in subtitle.text_frame.paragraphs:
    for run in p.runs:
        run.font.name = 'Microsoft YaHei'
        run.font.size = Pt(24)

# Slide 2: The Problem
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
set_title(slide, "行业痛点：当前中国白牌出海的“三大割裂死局”")

content2 = """
- 【洞察与生产割裂】(Information Gap)
  - 厂长懂开模，但不懂欧美年轻人的亚文化与审美变迁。
  - 亚马逊上的欧美差评痛点，无法有效穿透回国内的打样生产车间。
- 【账本与流量割裂】(Profit Gap)
  - 找货只看 1688 出厂价便宜，彻底忽视跨境极细的隐藏成本。
  - 被北美女装高达 25% 的退货率以及 FBA 抛货重税轻易拖垮利润底线。
- 【现存 AI 只是“辅助玩具”】(AI Implementation Gap)
  - 现有的“出海 AI”多停留在文案代写、机器翻译的单点工具层面。
  - 缺乏业务闭环能力，AI 无法下划触碰最硬核的财务风控与供应商谈判防坑。
"""
set_content(slide, content2)

# Slide 3: The Solution
slide = prs.slides.add_slide(slide_layout)
set_title(slide, "破局之道：真正的“端到端全域协同 (End-to-End Synergy)”")

content3 = """
- 从“单点聊天工具”进化为跨越业务壁垒的 Agentic 生意流
  - 摒弃“各管一段、写写文案”的弱智表层能力。
  - 本项目通过深度联动的多 Agent 网络，真实复现一家亿级出海大卖的内部核心运作：选品 -> 精算 -> 验厂 -> 包装。
- 数据流的高度穿透传递
  - 突破性实现不同 Agent 的能力传导闭环。
  - 前面做竞品差评大数据的 Agent 不止是产出报告，它的结论会无损传导给后面管供应链的 Agent，自动变成压迫 1688 代工厂“必须更换拉链工艺”的硬核验厂话术。
"""
set_content(slide, content3)

# Slide 4: Core Architecture
slide = prs.slides.add_slide(slide_layout)
set_title(slide, "CBEC Agent Matrix：四核数字高管")

content4 = """
- 1. 商业探测大脑 (Opportunity Spotter)
  - 抛除无意义的大盘，切入最核心的利润空间探测，锁定 $45-$65 等“价格真空带”。
- 2. 算账管家与死线 (ROI Assessor)
  - 重构真实出海跨境财务底账，将亚马逊重税与头尾程运费摊毁模型内嵌，划定生死的最高起订量与合理指导售价。
- 3. 美学映射图纸 (Community Mapper)
  - 打破族裔刻板印象，利用静奢风等普适产品特点，将工厂货精准发配至最高净值受众、并出具全盘出海架构指导。
- 4. 供应链风控猎手 (Sourcing Auditor)
  - 终极进货护城河：反骗话术、防伪雷达与阶梯式三轮实地打样 (PPS) QC 验审规则模板。
"""
set_content(slide, content4)

# Slide 5: The Workflow
slide = prs.slides.add_slide(slide_layout)
set_title(slide, "大厂级别的全漏斗闭环 (The Data-Driven Flywheel)")

content5 = """
- 以女装裙子出海（以“静奢结构裙”为例）的全链路智能飞轮：
  - 【输入】：100 条美区对真丝连衣裙的中差评分析数据
  - ⬇️ 【提取转化】Agent 诊断重灾区为：腹部剪裁紧绷、廉价起静电。
  - ⬇️ 【成本核算】Agent 精算空运体系及 FBA 大件重货惩罚红线标准。
  - ⬇️ 【倒逼验厂】Agent 向 1688 代工厂甩出极为冷酷的强制合规邮件：“要求把版型全部改做北美高标斜裁工艺且体积严控 2CM 厚度防抛”。
  - ⬇️ 【最终包装】生成具有极高情绪溢价与本地化美学质感的英文名字与社交名片定位。
"""
set_content(slide, content5)

# Slide 6: Moat & Value
slide = prs.slides.add_slide(slide_layout)
set_title(slide, "技术护城河与宏大商业价值")

content6 = """
- 极客级的底层 Guardrails (纪律栅栏) 防止大模型幻觉
  - 每个商业提议强制要求附带欧美市场第三方权威数据点作为交叉核验。
- 支持无极延展开源 `npx skills` API 连接全网能力
  - 可无缝挂载社区 3.2K 使用量的底层爬虫引擎及 SEO audit 组件，建立完整的商业基础设施。
- 斩断微笑曲线底端的枷锁 (最终目标)
  - 赋能千千万万只懂低头制造的中国厂长，利用 AI 拉平日渐消失的信息差。
  - 从靠内卷拼低价流血出海，彻底转变为手握定价话语权、拥有真实文化共鸣的全球化高质量 DTC 出海品牌矩阵集群！
"""
set_content(slide, content6)

# Save the presentation
prs.save('/Users/mamengyang/.gemini/antigravity/my project/阿里云创客大赛-跨境电商Agent操盘系统路演.pptx')
print("Presentation successfully generated!")
